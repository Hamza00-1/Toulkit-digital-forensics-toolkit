from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create a new presentation
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Aegis Digital Forensics Toolkit"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle.text = "An Enterprise-Grade, Dual-Architecture Triage Suite\n\nDepartment of Cybersecurity & Software Engineering"

    # Define a helper function for standard content slides
    def add_content_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(22)
            p.space_after = Pt(14)
            
        return slide

    # Slide 2: The Challenge
    add_content_slide(
        "The Challenge in Modern Forensics",
        [
            "Scale & Noise: Investigators must manually parse thousands of log lines and hidden metadata tags.",
            "Data Mutilation: Suspects format drives, delete files, and obfuscate payloads to hide tracks.",
            "Collaboration Bottlenecks: Heavy offline processing environments make remote reporting difficult.",
            "Our Goal: A solution that balances heavy offline processing with agile, secure reporting."
        ]
    )

    # Slide 3: Architecture
    add_content_slide(
        "Aegis Architecture: Desktop vs. Web",
        [
            "1. Python Desktop Client (Offline)",
            "  - Built with CustomTkinter for a modern, responsive feel.",
            "  - Deep OS integration ensures evidence never leaks to the internet.",
            "2. Flask Web Dashboard (Remote)",
            "  - Accessible via browser for non-technical managers.",
            "  - Secure RESTful API structure mirroring the desktop engine."
        ]
    )

    # Slide 4: Integrity
    add_content_slide(
        "Integrity & Chain of Custody",
        [
            "Principle: If evidence is altered without a log, it is inadmissible.",
            "Activity Tracker Database:",
            "  - Every forensic action triggers an automated INSERT into a secure MySQL environment.",
            "  - Captures timestamps, target files, success states, and the specific module used.",
            "  - Creates an immutable, automated audit trail for the investigator."
        ]
    )

    # Slide 5: Mathematical Validation
    add_content_slide(
        "Validation Over Trust: Hash & Threat Intel",
        [
            "The Cryptographic Validator Module",
            "1. Mathematical Proof",
            "  - Utilizes the hashlib library to compute MD5, SHA-1, and SHA-256 signatures.",
            "  - Proves file integrity before and after processing.",
            "2. Global Threat Intelligence",
            "  - Integrates directly with the VirusTotal API.",
            "  - We query global heuristic databases to flag known malware instantly, rather than relying on local assumptions."
        ]
    )

    # Slide 6: Automation
    add_content_slide(
        "Automating the Triage: Syslog Abstractor",
        [
            "Manual log analysis is prone to human error. Aegis automates this:",
            "  - Regex Pattern Matching: Extracts timestamps, hostnames, and messages from standard auth.log formats.",
            "  - Data Structuring: Maps unstructured text into powerful Pandas DataFrames.",
            "  - Anomaly Detection: Mathematically flags IPs with multiple failed SSH login attempts, highlighting brute-force vectors."
        ]
    )

    # Slide 7: Messy Reality
    add_content_slide(
        "Navigating the Messy Reality",
        [
            "Digital evidence is rarely pristine. Aegis recovers damaged data:",
            "  - Hexadecimal Data Carving: Bypasses corrupted file systems entirely. Scans raw memory dumps and byte-streams to isolate files using explicitly defined Magic Headers.",
            "  - Payload Decoder Matrix: Translates suspicious Base64, Hexadecimal, and Binary strings back into human-readable text using binascii."
        ]
    )

    # Slide 8: Live Demo
    add_content_slide(
        "Live Demonstration",
        [
            "1. Metadata Extraction (EXIF Analysis)",
            "2. Cryptographic Validation (EICAR Malware String Detection)",
            "3. Automated Log Abstraction (Syslog Analysis)",
            "4. Global Executive Dashboard & Web Application"
        ]
    )
    
    # Save the presentation
    file_name = 'Aegis_Forensics_Presentation.pptx'
    prs.save(file_name)
    print(f"Presentation saved successfully as {file_name}")

if __name__ == '__main__':
    create_presentation()
